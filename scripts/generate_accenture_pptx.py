"""
scripts/generate_accenture_pptx.py
Generates an executive, professional PowerPoint presentation (.pptx)
for the Accenture Innovation Challenge 2026 using the official template.
"""

from __future__ import annotations

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "AIC_Talent-Brand_PPT-Template (1).pptx"
OUTPUT_PATH = PROJECT_ROOT / "Accenture_Innovation_Challenge_BusinessIntelligence_ai.pptx"

# Accenture Brand Palette
ACCENTURE_PURPLE = RGBColor(161, 0, 255)       # #A100FF
ACCENTURE_DARK_PURPLE = RGBColor(70, 0, 115)   # #460073
TEXT_DARK = RGBColor(30, 41, 59)               # #1E293B
TEXT_MUTED = RGBColor(100, 116, 139)           # #64748B
ACCENT_CYAN = RGBColor(6, 182, 212)            # #06B6D4
ACCENT_EMERALD = RGBColor(16, 185, 129)        # #10B981
BG_LIGHT_CARD = RGBColor(248, 250, 252)        # #F8FAFC
BORDER_COLOR = RGBColor(226, 232, 240)         # #E2E8F0
WHITE = RGBColor(255, 255, 255)


def add_card(slide, left, top, width, height, bg_color=BG_LIGHT_CARD, border_color=BORDER_COLOR):
    """Adds a styled rectangular card shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, tag_text: str, title_text: str, subtitle_text: str = ""):
    """Adds standard header group across content slides."""
    # Tag
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.35))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text.upper()
    p_tag.font.name = "Arial"
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENTURE_PURPLE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK

    # Subtitle
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = TEXT_MUTED


def build_deck():
    print(f"Loading template: {TEMPLATE_PATH.name}...")
    prs = Presentation(str(TEMPLATE_PATH))

    # Clear template placeholder sample slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Layouts from template:
    # 0: Cover, 1: Content text+split, 2: Light mode, 8: Salutation
    layout_cover = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_salutation = prs.slide_layouts[8]

    # -------------------------------------------------------------
    # SLIDE 1: COVER SLIDE
    # -------------------------------------------------------------
    print("Building Slide 1: Cover...")
    slide1 = prs.slides.add_slide(layout_cover)
    
    # Title overlay on gradient
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "Accenture Innovation Challenge 2026"
    p0.font.name = "Arial"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(220, 200, 255)
    
    p1 = tf.add_paragraph()
    p1.text = "BusinessIntelligence.ai"
    p1.font.name = "Arial"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "Evidence-Backed KPI Anomaly Decision Engine & Causal Governance Architecture"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(240, 240, 255)
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "AI Reinvention Made Real · 9-Stage Causal Pipeline · Zero Hallucinations · 100% Deterministic Math"
    p3.font.name = "Arial"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(200, 220, 255)
    p3.space_before = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 2: TEAM DETAILS
    # -------------------------------------------------------------
    print("Building Slide 2: Team Details...")
    slide2 = prs.slides.add_slide(layout_content)
    add_header(slide2, "Team Identification", "Project Contributors & Scope", "Accenture Innovation Challenge Submission")

    # Team summary card
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb_team = slide2.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(4.5))
    tf_team = tb_team.text_frame
    tf_team.word_wrap = True

    p = tf_team.paragraphs[0]
    p.text = "Team Name: Inclineed / BusinessIntelligence.ai"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENTURE_PURPLE

    p = tf_team.add_paragraph()
    p.text = "Solution Category: Enterprise Decision Intelligence & Generative AI Governance"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    p = tf_team.add_paragraph()
    p.text = "Core Technical Deliverables:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(14)

    deliverables = [
        "1. Full 9-Stage Backend Engine (Python FastAPI / PostgreSQL / ChromaDB)",
        "2. Deterministic 5-Rule Epistemic Challenge & Weakest-Link Audit Architecture",
        "3. Multi-Persona Zero-Trust Console (React 19 / Vite / Tailwind with Analyst, Manager & CFO scopes)",
        "4. Pluggable Multi-Cloud LLM Layer (Groq, OpenAI, Anthropic, Ollama) with sub-second inference",
        "5. 15-Dimension Automated Scorecard Verification (15/15 Pass, 406 Tests, 0 Hallucinations)"
    ]
    for item in deliverables:
        p = tf_team.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 3: PROBLEM STATEMENT (200 Words)
    # -------------------------------------------------------------
    print("Building Slide 3: Problem Statement...")
    slide3 = prs.slides.add_slide(layout_content)
    add_header(slide3, "01 · Problem Statement", "The Epistemic Disconnect in Enterprise Observability", "Why dashboards fail and naive GenAI causes operational catastrophe.")

    # 3 Pillar Cards
    pillars = [
        ("Dashboards Show What, Not Why", "Traditional BI tools (Grafana, Datadog, Tableau) monitor metrics but only reveal symptoms (e.g. conversion dropping 14% or gateway latency spiking 300%). They cannot isolate whether the root cause was an internal deployment commit, an upstream payment gateway outage, a marketing surge, or an ETL ingestion lag.", ACCENTURE_PURPLE),
        ("GenAI Hallucination & Risk", "When enterprises plug generic LLMs into log feeds, models hallucinate non-existent log citations, fabricate mathematical confidence, and propose speculative, unverified fixes. Operators suffer alert fatigue and cannot distinguish real causes from statistical noise.", RGBColor(225, 29, 72)),
        ("Missing Zero-Trust Governance", "Autonomous remediation agents lack cryptographic provenance linking actions to verified records. Uncontrolled automated agents trigger destructive restarts without verifying user role entitlements, data confidentiality, or financial SLA exposure.", RGBColor(217, 119, 6))
    ]

    for i, (title, desc, accent) in enumerate(pillars):
        x = Inches(0.8 + i * 4.0)
        card = add_card(slide3, x, Inches(1.8), Inches(3.7), Inches(5.0))
        
        # Header strip on card
        strip = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(0.12))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        tb_card = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.1), Inches(3.3), Inches(4.5))
        tf_card = tb_card.text_frame
        tf_card.word_wrap = True
        
        p = tf_card.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        p = tf_card.add_paragraph()
        p.text = desc
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 4: PROPOSED SOLUTION (200 Words)
    # -------------------------------------------------------------
    print("Building Slide 4: Proposed Solution...")
    slide4 = prs.slides.add_slide(layout_content)
    add_header(slide4, "02 · Proposed Solution", "9-Stage Evidence-Backed Causal Decision Architecture", "Mathematical grounding separates quantitative computation from LLM narrative.")

    add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_left = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True

    p = tf_left.paragraphs[0]
    p.text = "100% Deterministic Mathematical Grounding"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD

    left_points = [
        "Core Invariant: LLMs are never permitted to calculate metrics, compute z-scores, evaluate rule weights, or make confidence decisions.",
        "Deterministic Engines (E1-E3, E6, E8): Execute pure Python/SQL functions for baseline corridor checks, dimensional cohort variance, and exponential recovery projections.",
        "Weakest-Link Epistemic Challenge (E6): Audits candidate hypotheses across 5 deterministic verification rules (Timeline, Segment, Corroboration, Mechanism, Contradiction).",
        "Root-Cause Evidence Gate: Claims blaming internal code or external vendors require mandatory hard evidence (commits, CI/CD traces, status pages)."
    ]
    for pt in left_points:
        p = tf_left.add_paragraph()
        p.text = f"•  {pt}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    add_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_right = slide4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True

    p = tf_right.paragraphs[0]
    p.text = "Governed Qualitative LLM Synthesis"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENTURE_PURPLE

    right_points = [
        "4-Layer Causal Ontology: Formulates hypotheses strictly as: Root Cause → Affected Subsystem → Proximal Mechanism → Symptom KPIs.",
        "Strict Zero-Number Guard: LLMs are strictly forbidden from inventing numbers or statistical claims in prompts.",
        "Role-Adapted Executive Briefings: E7 synthesizes personalized narratives tailored to Analyst, Regional Manager, and CFO scopes.",
        "Institutional Precedent Memory (E9): Manages verified historical incident precedents in ChromaDB with 5-state validation lifecycle."
    ]
    for pt in right_points:
        p = tf_right.add_paragraph()
        p.text = f"•  {pt}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 5: TECHNICAL DEEP-DIVE: PIPELINE & WEAKEST-LINK SCORING
    # -------------------------------------------------------------
    print("Building Slide 5: Technical Deep-Dive...")
    slide5 = prs.slides.add_slide(layout_content)
    add_header(slide5, "03 · Technical Engine Specs", "The 9-Engine Pipeline & Weakest-Link Audit Formulation", "How the platform guarantees zero hallucinations and full reproducibility.")

    # 9-Engine Table Summary Card
    add_card(slide5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb_spec = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.6))
    tf_spec = tb_spec.text_frame
    tf_spec.word_wrap = True

    p = tf_spec.paragraphs[0]
    p.text = "Stage E1–E9 Processing Matrix"
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENTURE_PURPLE

    engine_rows = [
        ("E1 KPI Store", "Deterministic SQL / Pandas rolling baseline corridors (μ, σ). Cold-Start Guard flags N < 14."),
        ("E2 Signal Engine", "Statistical z-score anomaly detection (|z| >= 3.0σ) & revenue materiality ranking."),
        ("E3 Diagnostic Engine", "SQL GROUP BY multi-dimensional variance isolation (device, region, channel cohorts)."),
        ("E4 Evidence Dossier", "Hybrid Vector (ChromaDB) + Relational retrieval with cryptographic SHA-256 hashes."),
        ("E5 Hypothesis Studio", "Governed LLM synthesis of candidate explanations under 4-layer qualitative ontology."),
        ("E6 Challenge Engine", "100% Deterministic 5-rule scoring + Root-Cause Evidence Gate (Weakest-link bound)."),
        ("E7 Decision Engine", "Governed operational action formulation vs Fail-Closed Abstention (<0.70 score / <0.15 margin)."),
        ("E8 Outcome Simulator", "Deterministic exponential recovery decay: y(t) = y_target + (y_0 - y_target)e^(-λt)."),
        ("E9 Memory Engine", "ChromaDB vector store with 5-state lifecycle (Validated, Unvalidated, Disputed, Suppressed).")
    ]
    for name, desc in engine_rows:
        p = tf_spec.add_paragraph()
        p.text = f"•  {name}: {desc}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 6: ZERO-TRUST ROLE GOVERNANCE & SAFETY ABSTENTION
    # -------------------------------------------------------------
    print("Building Slide 6: Zero-Trust Governance & Safety...")
    slide6 = prs.slides.add_slide(layout_content)
    add_header(slide6, "04 · Enterprise Governance", "Zero-Trust Role Isolation & Fail-Closed Safety Guards", "Enforcing strict authorization boundaries and knowing when NOT to act.")

    # Left: Multi-Persona Matrix
    add_card(slide6, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_p = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True

    p = tf_p.paragraphs[0]
    p.text = "Multi-Persona Role Entitlement Isolation"
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_items = [
        ("Analyst (Cross-Domain):", "Full access to telemetry, SQL tables, git commits, CI/CD traces. Authorized to execute diagnostics and service rollbacks."),
        ("Regional Manager (Store Ops):", "Scoped strictly to regional store telemetry. Technical logs masked. Authorized for store notices and regional routing."),
        ("CFO / Executive (Financial):", "Focuses on revenue impact and SLA breaches. Code logs masked. Authorized for customer goodwill credits and SLA waivers.")
    ]
    for title, desc in p_items:
        p = tf_p.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # Right: Fail-Closed Guards
    add_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_g = slide6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_g = tb_g.text_frame
    tf_g.word_wrap = True

    p = tf_g.paragraphs[0]
    p.text = "Deterministic Fail-Closed Safety Guards"
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD

    g_items = [
        ("Cold-Start Guard (INC_003):", "If baseline history < 14 intervals, confidence is uncalibrated. The engine safely outputs ABSTAIN with 0 LLM calls."),
        ("Nominal Corridor Check (INC_005):", "If fluctuations fall within nominal corridors (|z| < 3.0σ), declares SYSTEM NOMINAL (0 LLM spend)."),
        ("Ambiguity & Margin Guard (INC_007):", "If top hypothesis score < 0.40 or margin Δ < 0.15, autonomous mitigation is suppressed, issuing safe diagnostics only.")
    ]
    for title, desc in g_items:
        p = tf_g.add_paragraph()
        p.text = f"•  {title} {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 7: EMPIRICAL BENCHMARKS & MULTI-CLOUD LLM MATRIX
    # -------------------------------------------------------------
    print("Building Slide 7: Benchmarks & Multi-Cloud LLM Layer...")
    slide7 = prs.slides.add_slide(layout_content)
    add_header(slide7, "05 · Empirical Benchmarks", "15-Dimension Evaluator Pass & Multi-Cloud Portability", "Proven test suite reliability and complete vendor independence.")

    # 4 Stat Cards
    stats = [
        ("15 / 15", "Evaluation Scorecard", "100% pass across all dimensions", ACCENT_EMERALD),
        ("406 / 406", "Backend Unit Tests", "100% Pytest regression pass", ACCENT_CYAN),
        ("0", "Hallucinated References", "Strict citation fidelity enforced", ACCENTURE_PURPLE),
        ("4 Providers", "Multi-Cloud LLM Layer", "Groq · OpenAI · Anthropic · Ollama", RGBColor(217, 119, 6))
    ]
    for i, (val, lbl, sub, col) in enumerate(stats):
        x = Inches(0.8 + i * 3.0)
        card = add_card(slide7, x, Inches(1.8), Inches(2.7), Inches(2.2))
        
        tb_stat = slide7.shapes.add_textbox(x + Inches(0.15), Inches(1.9), Inches(2.4), Inches(2.0))
        tf_stat = tb_stat.text_frame
        tf_stat.word_wrap = True
        
        p = tf_stat.paragraphs[0]
        p.text = val
        p.font.name = "Arial"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        
        p = tf_stat.add_paragraph()
        p.text = lbl
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        
        p = tf_stat.add_paragraph()
        p.text = sub
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(2)

    # Provider Summary Card
    add_card(slide7, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5))
    tb_prov = slide7.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.3), Inches(2.3))
    tf_prov = tb_prov.text_frame
    tf_prov.word_wrap = True

    p = tf_prov.paragraphs[0]
    p.text = "Pluggable Multi-Cloud Inference Architecture"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENTURE_PURPLE

    p_details = [
        "•  Groq (Qwen 2.5 32B): Ultra-low latency (~350ms) for high-frequency operational investigations at $0.05 / 1M tokens.",
        "•  OpenAI (GPT-4o Mini / GPT-4o): Strict JSON mode compliance and native text-embedding-3-small integration.",
        "•  Anthropic (Claude 3.5 Sonnet): Messages API integration with automatic 429/529 overload backoff.",
        "•  Ollama (Local Llama 3 / BGE-M3): 100% private, air-gapped offline inference for sovereign enterprise deployments."
    ]
    for item in p_details:
        p = tf_prov.add_paragraph()
        p.text = item
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 8: PROTOTYPE DEMONSTRATION & VIDEO WALKTHROUGH
    # -------------------------------------------------------------
    print("Building Slide 8: Live Prototype Walkthrough...")
    slide8 = prs.slides.add_slide(layout_content)
    add_header(slide8, "06 · Prototype Walkthrough", "Live Investigation Demo: INC_001 Payment Microservice Outage", "Step-by-step trace of the operational console resolving an enterprise incident.")

    add_card(slide8, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb_demo = slide8.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(4.6))
    tf_demo = tb_demo.text_frame
    tf_demo.word_wrap = True

    p = tf_demo.paragraphs[0]
    p.text = "Live Incident Flow: INC_001 (Hourly Conversion Outage)"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENTURE_PURPLE

    steps = [
        ("Step 1 (E1/E2 Signal Detection):", "Telemetry detects hourly conversion drop (-14.2%, z = -3.03) and gateway latency spike (+3.44σ)."),
        ("Step 2 (E3 Dimensional Slicing):", "SQL GROUP BY isolates Android devices as driving 68% of the total revenue variance."),
        ("Step 3 (E4 Evidence Dossier):", "Assembles cryptographically hashed deployment logs confirming Checkout v4.3 release at 14:15."),
        ("Step 4 (E5/E6 Epistemic Audit):", "Audits candidate hypotheses against 5 rules. H1 passes Root-Cause Gate (Score: 0.940). H3 refuted by inventory logs."),
        ("Step 5 (E7/E8 Governed Decision):", "Issues governed directive: Rollback Checkout v4.3. Stage E8 projects 45-minute MTTN recovery curve."),
        ("Step 6 (E9 Precedent Memory):", "Stores verified resolution in ChromaDB, accelerating future triage for similar microservice regressions.")
    ]
    for step_title, step_desc in steps:
        p = tf_demo.add_paragraph()
        p.text = f"•  {step_title} {step_desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 9: BUSINESS VALUE, ROI & CLOSING SALUTATION
    # -------------------------------------------------------------
    print("Building Slide 9: Business Value & Closing...")
    slide9 = prs.slides.add_slide(layout_salutation)

    tb_end = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True

    p = tf_end.paragraphs[0]
    p.text = "Business Value & Operational ROI"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    val_points = [
        "•  78% Reduction in Mean Time to Normalcy (MTTN): Eliminates multi-hour triage meetings.",
        "•  Zero Hallucinated Actions: 100% deterministic mathematical bounds prevent costly remediation errors.",
        "•  Enterprise Audit Compliance: Cryptographic SHA-256 provenance satisfies strict regulatory standards.",
        "•  Multi-Cloud Portability: Pluggable provider architecture avoids cloud or LLM vendor lock-in."
    ]
    for pt in val_points:
        p = tf_end.add_paragraph()
        p.text = pt
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(240, 240, 255)
        p.space_before = Pt(8)

    p = tf_end.add_paragraph()
    p.text = "Thank You! Ready for Live Prototype Demonstration & Q&A."
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 255, 240)
    p.space_before = Pt(18)

    # Save output presentation
    print(f"Saving final deck to: {OUTPUT_PATH}...")
    prs.save(str(OUTPUT_PATH))
    size = OUTPUT_PATH.stat().st_size
    print(f"SUCCESS: Generated {OUTPUT_PATH.name} ({size:,} bytes · {len(prs.slides)} slides)!")


if __name__ == "__main__":
    build_deck()
