import sys
from pptx import Presentation

# 1. Inspect The_Decision_Pipeline_(3).pptx
old_ppt = Presentation(r"e:\accenture\The_Decision_Pipeline_(3).pptx")
print(f"=== The_Decision_Pipeline_(3).pptx has {len(old_ppt.slides)} slides ===")
for i, slide in enumerate(old_ppt.slides):
    title = slide.shapes.title.text if slide.shapes.title else "No Title"
    texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
    print(f"\n--- Slide {i+1}: Title: '{title}' ---")
    for t in texts:
        lines = t.strip().split("\n")
        for line in lines[:3]:
            print(f"   • {line[:100]}")

# 2. Inspect AIC_Talent-Brand_PPT-Template (1).pptx
template_ppt = Presentation(r"e:\accenture\AIC_Talent-Brand_PPT-Template (1).pptx")
print(f"\n=== AIC_Talent-Brand_PPT-Template (1).pptx has {len(template_ppt.slides)} sample slides and {len(template_ppt.slide_layouts)} layouts ===")
for i, slide in enumerate(template_ppt.slides):
    title = slide.shapes.title.text if slide.shapes.title else "No Title"
    print(f"Template Slide {i+1}: Title: '{title}'")
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            print(f"   [shape] {shape.text.strip()[:80]}")

print("\nAvailable Layouts in Template:")
for i, layout in enumerate(template_ppt.slide_layouts):
    print(f"Layout {i}: '{layout.name}'")
